"""Apply HANDOFF fixes to slides/Math130Unit3D.pptx.

Run from repo root:
    python3 slides/_math130_unit3d_wip/build_pptx.py

Idempotent re-runs are supported: the script always opens the original
and rebuilds. The ORIGINAL is read from the .orig snapshot so repeat
runs produce the same output.
"""
from __future__ import annotations

import os
import shutil
from copy import deepcopy
from lxml import etree

from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor


REPO = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
SRC  = os.path.join(REPO, "slides", "Math130Unit3D.pptx")
ORIG = os.path.join(REPO, "slides", "_math130_unit3d_wip", "Math130Unit3D.orig.pptx")
IMG  = os.path.join(REPO, "slides", "_math130_unit3d_wip", "img")


# Palette
NAVY   = RGBColor(0x1E, 0x3A, 0x5F)
TEAL   = RGBColor(0x08, 0x91, 0xB2)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
RED    = RGBColor(0xDC, 0x26, 0x26)
MUTED  = RGBColor(0x64, 0x74, 0x8B)
BG     = RGBColor(0xF8, 0xFA, 0xFC)
DARK   = RGBColor(0x0F, 0x2B, 0x46)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


def replace_picture(slide, pic_shape_idx: int, image_path: str):
    """Remove picture at index, insert a fresh one at the same rect."""
    shape = slide.shapes[pic_shape_idx]
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes.add_picture(image_path, left, top, width, height)


def set_text(shape, text: str, *, size=None, color=None,
             bold=None, italic=None, font_name=None,
             align=None, preserve_empty=True):
    """Replace text in a shape. Supports multi-line via `|` or newlines."""
    tf = shape.text_frame
    # Preserve first paragraph's alignment by default
    tf.clear()
    lines = text.split("\n") if "\n" in text else text.split(" | ")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align is not None:
            p.alignment = align
        if not line and preserve_empty:
            continue
        r = p.add_run()
        r.text = line
        if size is not None:  r.font.size = Pt(size)
        if color is not None: r.font.color.rgb = color
        if bold is not None:  r.font.bold = bold
        if italic is not None: r.font.italic = italic
        if font_name is not None: r.font.name = font_name


def style_runs(shape, *, color=None, bold=None, italic=None, size=None):
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if color is not None: r.font.color.rgb = color
            if bold is not None:  r.font.bold = bold
            if italic is not None: r.font.italic = italic
            if size is not None:  r.font.size = Pt(size)


def reorder_slides(prs, new_order):
    """Reorder slides per a 0-indexed permutation of current positions."""
    sldIdLst = prs.slides._sldIdLst
    children = list(sldIdLst)
    assert sorted(new_order) == list(range(len(children))), "bad perm"
    # Detach all
    for c in children:
        sldIdLst.remove(c)
    for idx in new_order:
        sldIdLst.append(children[idx])


def delete_shape(shape):
    sp = shape._element
    sp.getparent().remove(sp)


def recolor_run(shape, needle: str, new_color: RGBColor):
    """Recolor every run in `shape` whose text contains `needle`."""
    if not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if needle in r.text:
                r.font.color.rgb = new_color


# ---------------------------------------------------------------------------
# Load from original snapshot so re-runs are deterministic.
# ---------------------------------------------------------------------------
if not os.path.exists(ORIG):
    shutil.copy(SRC, ORIG)
prs = Presentation(ORIG)

# Indexed references (all 0-indexed on the CURRENT 13-slide deck).
S_TITLE    = 0   # Title
S_INTRO    = 1   # Section 8.4 / Vectors (will become "Why Vectors?")
S_SKILLS   = 2   # Skills Overview
S_DEFS     = 3   # Vectors and Scalars
S_COMP5    = 4   # Writing a Vector in Component Form
S_FINDCMP  = 5   # Finding Components from a Graph  (repurpose -> Components from |v|,θ)
S_ADD7     = 6   # Vector Addition & Scalar Multiplication
S_MAG8     = 7   # Magnitude and Direction
S_QUAD9    = 8   # Direction Angle - Quadrant Adjustment
S_DIRANG   = 9   # Direction Angle of ai + bj
S_SCAL11   = 10  # Scalar Multiplication - Visual
S_TESTPREP = 11  # Test Prep
S_SUMMARY  = 12  # Module 13 Summary


# ---------------------------------------------------------------------------
# 1. Swap the five embedded diagrams to the corrected PNGs.
#    (Slide 5's swap happens LATER, after we clone it for "You Try".)
# ---------------------------------------------------------------------------
replace_picture(prs.slides[S_ADD7],  13, f"{IMG}/s_vector_addition.png")
replace_picture(prs.slides[S_MAG8],   3, f"{IMG}/s_magnitude_direction.png")
replace_picture(prs.slides[S_QUAD9],  1, f"{IMG}/s_quadrant_arrows.png")
replace_picture(prs.slides[S_SCAL11], 1, f"{IMG}/s_scalar_multiplication.png")


# ---------------------------------------------------------------------------
# 1b. Slide 9 (Quadrant Adjustment) - replace the 4-example breakdown
#     with the single compact rule from HANDOFF §3.8.
#
#     After the picture swap above, shape indices are:
#       0 title | 1 right-panel bg | 2 "Examples" label
#       3..14   (four rows: <a,b>, Q label, θ value)
#       15      NEW PICTURE
# ---------------------------------------------------------------------------
s = prs.slides[S_QUAD9]
# Re-label the header (was "Examples")
set_text(s.shapes[2], "One Rule for θ", size=20, color=NAVY,
         bold=True, font_name="Georgia")
# Drop rows 3..14 (12 shapes) in reverse order — keep the picture at [15].
for idx in range(14, 2, -1):
    delete_shape(s.shapes[idx])
# Add one new text block with the rule.
rule_box = s.shapes.add_textbox(Emu(5212080), Emu(1463040),
                                Emu(3657600), Emu(2925000))
tf = rule_box.text_frame
tf.word_wrap = True
lines = [
    ("Compute α = tan⁻¹(b/a), then:", NAVY, 16, False),
    (" ", NAVY, 10, False),
    ("if a < 0  →  θ = α + 180°", PURPLE, 16, True),
    ("     (Q II or Q III)", MUTED, 13, False),
    (" ", NAVY, 10, False),
    ("else if b < 0  →  θ = α + 360°", PURPLE, 16, True),
    ("     (Q IV)", MUTED, 13, False),
    (" ", NAVY, 10, False),
    ("otherwise  →  θ = α", PURPLE, 16, True),
    ("     (Q I)", MUTED, 13, False),
]
for i, (text, color, size, bold) in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = "Calibri"


# ---------------------------------------------------------------------------
# 1c. Palette audit (HANDOFF §3.9):
#     TEAL = definitions/formulas, ORANGE = examples/worked problems.
#     The original deck has these swapped on slides 5 & 8 — fix them.
# ---------------------------------------------------------------------------
# Slide 5 (Writing a Vector in Component Form): "Formula:" -> TEAL, "Example:" -> ORANGE
recolor_run(prs.slides[S_COMP5].shapes[2], "Formula:", TEAL)
recolor_run(prs.slides[S_COMP5].shapes[5], "Example:", ORANGE)
# Slide 8 (Magnitude and Direction): "Magnitude:" already TEAL; recolor "Direction:"
#  from orange to teal (both are formulas), and promote "Example:" to ORANGE.
recolor_run(prs.slides[S_MAG8].shapes[2], "Direction:", TEAL)
recolor_run(prs.slides[S_MAG8].shapes[5], "Example:", ORANGE)


# ---------------------------------------------------------------------------
# 2. Slide 3 - Skills Overview: relabel skill 5.
# ---------------------------------------------------------------------------
s = prs.slides[S_SKILLS]
# Shape 15 is the skill-5 label.
set_text(s.shapes[15],
         "(Supplementary) Finding the components of a vector given magnitude and angle",
         size=16, color=NAVY, font_name="Calibri")


# ---------------------------------------------------------------------------
# 3. Slide 2 - repurpose the redundant "Section 8.4 / Vectors" card as
#    "Why Vectors?" motivation. Dark background -> use WHITE / TEAL text.
# ---------------------------------------------------------------------------
s = prs.slides[S_INTRO]
# shape 1 = "Section 8.4"   -> headline: "Why Vectors?"
# shape 2 = "Vectors"       -> one-line lede
# shape 3 = subtitle        -> bullet block
set_text(s.shapes[1], "Why Vectors?", size=44, color=WHITE, bold=True, font_name="Georgia")
set_text(s.shapes[2], "Many real-world quantities need both size AND direction.",
         size=22, color=TEAL, font_name="Calibri")
# Widen shape 3 so the bullets fit the whole column.
s.shapes[3].width  = Emu(7772400)
s.shapes[3].height = Emu(1828800)
set_text(
    s.shapes[3],
    "•  Velocity — wind at 15 mph to the north-east\n"
    "•  Force — 50 N pulling down at 30°\n"
    "•  Displacement — walking 3 blocks east, then 2 north\n"
    " \n"
    "Scalars (speed, mass, temperature) need only size.",
    size=18, color=WHITE, font_name="Calibri",
)


# ---------------------------------------------------------------------------
# 4. Slide 4 - Vectors and Scalars: append unit-vector + zero-vector info
#    inside the existing "Notation" card (shape 11).
# ---------------------------------------------------------------------------
s = prs.slides[S_DEFS]
notation = s.shapes[11]
set_text(
    notation,
    "Component form:  v = ⟨a, b⟩  =  a i + b j\n"
    "Unit vectors:  i = ⟨1, 0⟩,  j = ⟨0, 1⟩\n"
    "Zero vector:  0 = ⟨0, 0⟩  (no magnitude, no direction)\n"
    "Magnitude:  |v| = √(a² + b²)     Direction angle:  θ  (CCW from +x)",
    size=16, color=NAVY, font_name="Calibri",
)


# ---------------------------------------------------------------------------
# 5. Slide 6 - repurpose "Finding Components from a Graph" as
#    "Components from |v| and θ". Drop the redundant Method-1 column,
#    drop the stray θ ≈ 143.1°, drop double bullets, swap in the new image.
# ---------------------------------------------------------------------------
s = prs.slides[S_FINDCMP]
# Title
set_text(s.shapes[0], "Components from |v| and θ",
         size=32, color=NAVY, bold=True, font_name="Georgia")

# Left column (old Method 1): reuse the card as the new formula card.
set_text(s.shapes[3], "Given |v| and θ", size=20, color=TEAL,
         bold=True, font_name="Calibri")
set_text(s.shapes[4],
         "vₓ = |v| cos θ\n"
         "vᵧ = |v| sin θ\n"
         "v = ⟨|v| cos θ, |v| sin θ⟩",
         size=17, color=NAVY, font_name="Calibri")

# Right column (old Method 2): DELETE the four shapes outright -
# the diagram replaces them.
for idx in sorted([5, 6, 7, 8], reverse=True):
    delete_shape(s.shapes[idx])

# Bottom card: keep the example slot, replace the awful Method-1/Method-2
# example with a |v|,θ worked example and no theta line from components.
# After deleting 4 shapes above, indices shift: old 9,10,11 -> 5,6,7.
set_text(s.shapes[6], "Example:  |v| = 5,  θ = 53.1°",
         size=18, color=ORANGE, bold=True, font_name="Calibri")
set_text(s.shapes[7],
         "vₓ = 5 cos 53.1° ≈ 3\n"
         "vᵧ = 5 sin 53.1° ≈ 4\n"
         "v = ⟨3, 4⟩",
         size=16, color=NAVY, font_name="Calibri")

# Embed the new diagram in the area vacated by the "When to use" card.
# Use width-only so python-pptx preserves the image's native aspect.
# Slot: L=4.8..9.8in, T=1.0..3.4in  -> ~2.4in tall. Use width 2.6in centered.
prs.slides[S_FINDCMP].shapes.add_picture(
    f"{IMG}/s_components_from_mag_theta.png",
    Emu(5600000), Emu(914400), width=Emu(2468880),
)


# ---------------------------------------------------------------------------
# 6. Slide 13 (Summary) - the original "Next: Unit 3 Review -> Test 3" sat
#    BEHIND the "Key for Test 3" card. Move it DOWN below that card and
#    bump the contrast.
#
#    "Key for Test 3" card bottom edge: T=4023360 + H=685800 = 4709160
#    Slide height (5.625 in): 5143500. Footer height 365760 leaves room.
# ---------------------------------------------------------------------------
s = prs.slides[S_SUMMARY]
footer = s.shapes[16]
footer.top = Emu(4760000)
footer.height = Emu(365760)
# Recolor and bold for contrast (teal on dark navy background).
for p in footer.text_frame.paragraphs:
    for r in p.runs:
        r.font.color.rgb = TEAL
        r.font.bold = True
        r.font.italic = False
        if r.font.size is None:
            r.font.size = Pt(16)


# ---------------------------------------------------------------------------
# 7. Add a new "You Try" slide as a duplicate of the existing slide 5
#    (Component Form), which has the same title + side-image layout.
# ---------------------------------------------------------------------------
def clone_slide(prs, src_slide):
    """XML-level clone of a slide, appended at the end."""
    blank_layout = src_slide.slide_layout
    new_slide = prs.slides.add_slide(blank_layout)
    # Strip any placeholders introduced by the layout
    for shp in list(new_slide.shapes):
        sp = shp._element
        sp.getparent().remove(sp)
    # Copy shapes from source
    for shp in src_slide.shapes:
        new_slide.shapes._spTree.append(deepcopy(shp._element))
    return new_slide


you_try = clone_slide(prs, prs.slides[S_COMP5])
# you_try currently mirrors pristine slide 5 layout.
# Shape indices:
#   0 = title        1 = formula card bg       2 = formula text
#   3 = PICTURE      4 = example card bg       5 = example body text
set_text(you_try.shapes[0], "You Try",
         size=32, color=NAVY, bold=True, font_name="Georgia")

# Change the dark formula-card bg to a light card so NAVY text is readable.
prompt_card = you_try.shapes[1]
prompt_card.top    = Emu(914400)
prompt_card.height = Emu(914400)   # ~1.0 in - plenty for 2 lines
prompt_card.fill.solid()
prompt_card.fill.fore_color.rgb = WHITE
prompt_card.line.color.rgb = TEAL
# Keep prompt text to two lines at 15pt.
txt = you_try.shapes[2]
txt.top    = Emu(960120)
txt.height = Emu(822960)
txt.left   = Emu(640080)
txt.width  = Emu(4023360)
set_text(txt,
         "Vector v goes from P(−4, −2) to Q(2, 1).  Find:\n"
         "(a) component form   (b) magnitude   (c) direction angle",
         size=15, color=NAVY, font_name="Calibri")

set_text(you_try.shapes[5],
         "(a)  v = ⟨2 − (−4), 1 − (−2)⟩ = ⟨6, 3⟩\n"
         "(b)  |v| = √(36 + 9) = √45 = 3√5 ≈ 6.71\n"
         "(c)  θ = tan⁻¹(3 / 6) ≈ 26.57°   (Q I, no adjustment)",
         size=16, color=NAVY, font_name="Calibri")
# Picture swap happens LAST to keep text-shape indices stable.
replace_picture(you_try, 3, f"{IMG}/s_you_try.png")

# Now do the deferred slide-5 picture swap.
replace_picture(prs.slides[S_COMP5], 3, f"{IMG}/s_component_form.png")


# ---------------------------------------------------------------------------
# 8. Reorder slides to match HANDOFF's proposed 14-slide order.
#
#   new#  old-index  content
#     1     0   Title
#     2     1   Why Vectors?
#     3     2   Skills Overview
#     4     3   Definitions
#     5     4   Component Form  (img swapped)
#     6     6   Addition & Scalar Mult  (img swapped)
#     7    10   Scalar x Visual  (img swapped)
#     8     7   Magnitude & Direction  (img swapped)
#     9     8   Quadrant Adjustment  (img swapped)
#    10     9   Direction Angle of ai+bj
#    11     5   Components from |v|,θ  (repurposed)
#    12    13   You Try  (appended)
#    13    11   Test Prep
#    14    12   Summary
# ---------------------------------------------------------------------------
NEW_ORDER = [0, 1, 2, 3, 4, 6, 10, 7, 8, 9, 5, 13, 11, 12]
reorder_slides(prs, NEW_ORDER)


# ---------------------------------------------------------------------------
# Save.
# ---------------------------------------------------------------------------
prs.save(SRC)
print(f"Wrote {SRC}")
