"""Apply all HANDOFF fixes to slides/Math130Unit3C.pptx.

Run from repo root:
    python3 slides/_math130_unit3c_wip/build_pptx.py

Idempotent: always reads from the .orig snapshot so re-runs produce the same output.
"""
from __future__ import annotations

import os
import shutil
from copy import deepcopy

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

REPO = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
SRC  = os.path.join(REPO, "slides", "Math130Unit3C.pptx")
ORIG = os.path.join(REPO, "slides", "_math130_unit3c_wip", "Math130Unit3C.orig.pptx")
IMG  = os.path.join(REPO, "slides", "_math130_unit3c_wip", "img")

# Palette
NAVY   = RGBColor(0x1E, 0x3A, 0x5F)
TEAL   = RGBColor(0x08, 0x91, 0xB2)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
GREEN  = RGBColor(0x05, 0x96, 0x69)
RED    = RGBColor(0xDC, 0x26, 0x26)
RED_BG = RGBColor(0xFF, 0xED, 0xED)
MUTED  = RGBColor(0x64, 0x74, 0x8B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x0F, 0x2B, 0x46)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def set_text(shape, text: str, *, size=None, color=None,
             bold=None, italic=None, font_name=None, align=None):
    """Replace all text in a shape. Use '\\n' for line breaks."""
    tf = shape.text_frame
    tf.clear()
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align is not None:
            p.alignment = align
        r = p.add_run()
        r.text = line
        if size      is not None: r.font.size       = Pt(size)
        if color     is not None: r.font.color.rgb  = color
        if bold      is not None: r.font.bold        = bold
        if italic    is not None: r.font.italic      = italic
        if font_name is not None: r.font.name        = font_name


def set_runs(shape, segments: list[tuple[str, RGBColor | None, float | None, bool | None]]):
    """
    Replace text with multiple styled runs on successive paragraphs.
    segments: list of (text, color, size_pt, bold)
    A segment whose text is '' inserts a blank paragraph.
    """
    tf = shape.text_frame
    tf.clear()
    first = True
    for text, color, size_pt, bold in segments:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if not text:
            continue
        r = p.add_run()
        r.text = text
        if color   is not None: r.font.color.rgb = color
        if size_pt is not None: r.font.size      = Pt(size_pt)
        if bold    is not None: r.font.bold       = bold
        r.font.name = "Calibri"


def recolor_runs_containing(shape, needle: str, new_color: RGBColor):
    """Recolor every run whose text contains `needle`."""
    if not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if needle in r.text:
                r.font.color.rgb = new_color


def delete_shape(shape):
    shape._element.getparent().remove(shape._element)


def clone_slide(prs, src_slide):
    """Append an XML-level clone of src_slide at the end of prs."""
    new_slide = prs.slides.add_slide(src_slide.slide_layout)
    for shp in list(new_slide.shapes):
        shp._element.getparent().remove(shp._element)
    for shp in src_slide.shapes:
        new_slide.shapes._spTree.append(deepcopy(shp._element))
    return new_slide


def insert_slide_at(prs, slide, position: int):
    """Move a slide (already in prs) from its current position to `position` (0-indexed)."""
    sldIdLst = prs.slides._sldIdLst
    entries = list(sldIdLst)
    # Find the entry for this slide's rId
    slide_rId = None
    for entry in entries:
        rId = entry.get("r:id")
        if prs.slides._sldIdLst.getparent() is not None:
            pass
        # Match by checking the slide object
    # Simpler: rebuild list with desired order
    # Get all slide rIds in current order
    from lxml import etree
    ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    all_entries = list(sldIdLst)
    # The newly cloned slide is at the end
    entry_to_move = all_entries[-1]
    all_entries.remove(entry_to_move)
    all_entries.insert(position, entry_to_move)
    for e in list(sldIdLst):
        sldIdLst.remove(e)
    for e in all_entries:
        sldIdLst.append(e)


def reorder_slides(prs, new_order: list[int]):
    """Reorder slides by 0-indexed permutation of CURRENT positions."""
    sldIdLst = prs.slides._sldIdLst
    children = list(sldIdLst)
    assert sorted(new_order) == list(range(len(children))), \
        f"bad permutation: {new_order}"
    for c in children:
        sldIdLst.remove(c)
    for idx in new_order:
        sldIdLst.append(children[idx])


def add_textbox(slide, left_in, top_in, width_in, height_in):
    return slide.shapes.add_textbox(
        Emu(int(left_in * 914400)),
        Emu(int(top_in  * 914400)),
        Emu(int(width_in  * 914400)),
        Emu(int(height_in * 914400)),
    )


# ---------------------------------------------------------------------------
# Load from original snapshot so re-runs are deterministic.
# ---------------------------------------------------------------------------
if not os.path.exists(ORIG):
    shutil.copy(SRC, ORIG)

prs = Presentation(ORIG)

# Current 17-slide indices (0-based)
S_TITLE    = 0   # Module 12 title
S_INTRO    = 1   # Section 7.1 / 7 ALEKS Skills
S_SKILLS   = 2   # Skills Overview
S_STRATEGY = 3   # Problem-Solving Strategy
S_SIDE     = 4   # Finding a Side Length (Skill 1)
S_ANGLE    = 5   # Finding an Angle / Inverse Trig (Skill 4)
S_ELEV     = 6   # Angles of Elevation and Depression (concept)
S_DEPRESS  = 7   # Example — Angle of Depression
S_DRT      = 8   # Trig Functions and d = rt (concept)
S_DRTEX    = 9   # Example — Trig and d = rt
S_SOLV     = 10  # Solving a Right Triangle (concept)
S_SOLVEX   = 11  # Example — Solving a Right Triangle
S_TWO      = 12  # Word Problem — Two Right Triangles
S_TESTEX   = 13  # Test Prep — Exact Values
S_TESTSOLV = 14  # Test Prep — Solve a Right Triangle
S_TESTELEV = 15  # Test Prep — Elevation & Depression
S_SUMMARY  = 16  # Module 12 Summary


# ===========================================================================
# FIX 1: Slide 8 (Depression example) — transition line from MUTED to TEAL
# "Depression from top = elevation from ship = 15°" should be TEAL (key fact)
# ===========================================================================
s = prs.slides[S_DEPRESS]
recolor_runs_containing(s.shapes[5], "Depression from top", TEAL)


# ===========================================================================
# FIX 2: Slide 9 (d=rt concept) — harmonize step colors
# "Horizontal" and "Vertical" headers: both are formula-section headers → TEAL
# (Vertical was PURPLE, which is reserved for final results)
# ===========================================================================
s = prs.slides[S_DRT]
recolor_runs_containing(s.shapes[10], "Vertical", TEAL)


# ===========================================================================
# FIX 3: Slide 13 (Two Right Triangles) — insert missing algebra step
# Between "d(1.327) = d(0.601) + 50(0.601)" and "d(0.726) = 30.04"
# insert "d(1.327 − 0.601) = 50 × 0.601"
# ===========================================================================
s = prs.slides[S_TWO]
tb = s.shapes[5]  # main solution text block
tf = tb.text_frame

# Rebuild the paragraph list with the extra line inserted
OLD_LINES = [
    ("Set up two equations:", TEAL, 16.5, True),
    ("Closer: tan 53° = h / d", TEAL, 16.5, False),
    ("Farther: tan 31° = h / (d + 50)", ORANGE, 16.5, False),
    ("", None, 16.5, None),
    ("Solve: h = d tan 53° = (d+50) tan 31°", PURPLE, 16.5, True),
    ("d(1.327) = d(0.601) + 50(0.601)", None, 15, None),
    ("d(1.327 \u2212 0.601) = 50 \u00d7 0.601", MUTED, 14, False),   # NEW intermediate step
    ("d(0.726) = 30.04", None, 15, None),
    ("d \u2248 41.4 ft", None, 15, None),
    ("", None, 15, None),
    ("h = 41.4 \u00d7 tan 53\u00b0 \u2248 41.4 \u00d7 1.327", TEAL, 15, False),
    ("h \u2248 54.9 ft", ORANGE, 16.5, True),
]
set_runs(tb, OLD_LINES)


# ===========================================================================
# FIX 4: Slide 6 (Inverse Trig / Finding an Angle) — upgrade DEGREE warning
# Current: 13pt red text, easy to miss.
# New: larger callout box with a light-red background.
# ===========================================================================
s = prs.slides[S_ANGLE]
# The warning is the last part of shapes[5] text. Rebuild that text block
# keeping the worked solution, then add a separate styled warning box.
tb = s.shapes[5]
set_runs(tb, [
    ("sin A = opp / hyp = 7 / 15 \u2248 0.4667", None, 18, None),
    ("", None, 10, None),
    ("A = sin\u207b\u00b9(0.4667)", None, 18, None),
    ("A \u2248 27.8\u00b0", ORANGE, 20, True),
    ("", None, 10, None),
    ("Use sin\u207b\u00b9 (arcsin), cos\u207b\u00b9, or tan\u207b\u00b9 to find angles.", MUTED, 14, False),
])

# Add a warning callout box below the worked solution
warn_box = add_textbox(s, 0.7, 3.85, 5.0, 0.65)
warn_tf = warn_box.text_frame
warn_tf.word_wrap = True
# Light-red fill + red border
warn_box.fill.solid()
warn_box.fill.fore_color.rgb = RGBColor(0xFF, 0xED, 0xED)
warn_box.line.color.rgb = RED
warn_box.line.width = Emu(19050)  # 1.5pt
p = warn_tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "\u26a0  Calculator must be in DEGREE mode when using inverse trig!"
r.font.size = Pt(15)
r.font.bold = True
r.font.color.rgb = RED
r.font.name = "Calibri"


# ===========================================================================
# FIX 5: Slide 15 (Test Prep — Solve a Right Triangle)
# "Solution:" label color: change from ORANGE to GREEN (it's a final answer,
# not an example header)
# ===========================================================================
s = prs.slides[S_TESTSOLV]
recolor_runs_containing(s.shapes[4], "Solution:", GREEN)


# ===========================================================================
# FIX 6: Slide 17 (Summary) footer — make visible
# "Next: Module 13 — Section 8.4 (Vectors in Component Form)"
# ===========================================================================
s = prs.slides[S_SUMMARY]
footer = s.shapes[13]
footer.top    = Emu(int(4.82 * 914400))
footer.height = Emu(int(0.45 * 914400))
for p in footer.text_frame.paragraphs:
    for r in p.runs:
        r.font.color.rgb = TEAL
        r.font.bold      = True
        r.font.italic    = False
        if r.font.size is None or r.font.size < Pt(15):
            r.font.size = Pt(15)


# ===========================================================================
# NEW SLIDE A: Word Problem — One Right Triangle (Skill 2)
# Clone Slide 5 (Finding a Side) as the template — same two-column layout.
# Position: after Slide 5, before Slide 6 (becomes new index 5).
# ===========================================================================
skill2 = clone_slide(prs, prs.slides[S_SIDE])

# Title
set_text(skill2.shapes[0],
         "Word Problem \u2014 One Right Triangle",
         size=33, color=NAVY, bold=True, font_name="Georgia")

# Problem statement banner (shape[2] on S_SIDE = wide orange banner)
set_text(skill2.shapes[2],
         "A 20 ft ladder leans against a wall at 65\u00b0 from the ground. "
         "How high up the wall does it reach?",
         size=17, color=WHITE, bold=False, font_name="Calibri")

# Main solution block (shape[4])
set_runs(skill2.shapes[4], [
    ("Step 1: Identify", TEAL, 17, True),
    ("Hyp = 20 ft,  angle = 65\u00b0,  find opposite (wall height) \u2192 use sin (SOH)", None, 16, None),
    ("", None, 10, None),
    ("Step 2: Set up", ORANGE, 17, True),
    ("sin 65\u00b0 = h / 20", None, 17, None),
    ("", None, 10, None),
    ("Step 3: Solve", PURPLE, 17, True),
    ("h = 20 \u00d7 sin 65\u00b0 = 20 \u00d7 0.9063 \u2248 18.1 ft", None, 17, None),
    ("", None, 10, None),
    ("Check:", GREEN, 17, True),
    ("18.1 < 20 \u2713 (opposite < hypotenuse)", None, 15, False),
])

# Swap the embedded picture to the ladder diagram
# shape[3] on the cloned slide is the picture (same index as S_SIDE)
pic_shape = skill2.shapes[3]
left, top, width, height = (pic_shape.left, pic_shape.top,
                             pic_shape.width, pic_shape.height)
pic_shape._element.getparent().remove(pic_shape._element)
skill2.shapes.add_picture(
    os.path.join(IMG, "s_skill2_ladder.png"),
    left, top, width, height)

# Insert the cloned slide at position 5 (after old S_SIDE=4)
insert_slide_at(prs, skill2, 5)
# Slide indices have now shifted by +1 for slides after position 5.
# Adjust our constants:
S_ANGLE    += 1   # 6
S_ELEV     += 1   # 7
S_DEPRESS  += 1   # 8
S_DRT      += 1   # 9
S_DRTEX    += 1   # 10
S_SOLV     += 1   # 11
S_SOLVEX   += 1   # 12
S_TWO      += 1   # 13
S_TESTEX   += 1   # 14
S_TESTSOLV += 1   # 15
S_TESTELEV += 1   # 16
S_SUMMARY  += 1   # 17


# ===========================================================================
# NEW SLIDE B: You Try — angle-of-elevation practice problem
# Clone Slide 8 (Depression example) as the template — same layout with a
# picture on the right and a solution block on the left.
# Position: after Slide 13 (Two Right Triangles) → becomes new index 14.
# ===========================================================================
you_try = clone_slide(prs, prs.slides[S_TWO])

set_text(you_try.shapes[0],
         "You Try",
         size=36, color=NAVY, bold=True, font_name="Georgia")

set_text(you_try.shapes[2],
         "A surveyor stands 80 ft from a building. "
         "The angle of elevation to the roof is 42\u00b0. Find the building height.",
         size=16, color=WHITE, bold=False, font_name="Calibri")

set_runs(you_try.shapes[5], [
    ("Set up:", TEAL, 17, True),
    ("tan 42\u00b0 = h / 80", None, 17, None),
    ("", None, 10, None),
    ("Solve:", PURPLE, 17, True),
    ("h = 80 \u00d7 tan 42\u00b0", None, 17, None),
    ("h = 80 \u00d7 0.9004", None, 17, None),
    ("h \u2248 72.0 ft", ORANGE, 18, True),
    ("", None, 10, None),
    ("Check:", GREEN, 17, True),
    ("72.0 ft < 80 ft  \u2713  (height < base distance for 42\u00b0 < 45\u00b0)", None, 14, False),
])

# Replace the picture
pic = you_try.shapes[3]
left, top, width, height = pic.left, pic.top, pic.width, pic.height
pic._element.getparent().remove(pic._element)
you_try.shapes.add_picture(
    os.path.join(IMG, "s_you_try.png"),
    left, top, width, height)

# Insert after S_TWO (index 13), so position = 14
insert_slide_at(prs, you_try, 14)


# ===========================================================================
# Save
# ===========================================================================
prs.save(SRC)
print(f"Wrote {SRC}  ({len(prs.slides)} slides)")
